from __future__ import annotations
import os
import sys
import json
import uuid
import hashlib
import shutil
from pathlib import Path
from typing import Iterable, List, Optional, Dict, Any
import fitz  # PyMuPDF
import pandas as pd
import docx
from pptx import Presentation
from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from utils.model_loader import ModelLoader
from logger import GLOBAL_LOGGER as log
from exception.custom_exception import DocumentPortalException
from utils.file_io import generate_session_id, save_uploaded_files
from utils.document_ops import load_documents,load_documents_from_sql,load_excel_or_csv_as_documents

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".txt", ".pptx", ".csv", ".xlsx", ".xls", ".md"
}

# FAISS Manager (load-or-create)
class FaissManager:
    def __init__(self, index_dir: Path, model_loader: Optional[ModelLoader] = None):
        self.index_dir = Path(index_dir)
        self.index_dir.mkdir(parents=True, exist_ok=True)
        
        self.meta_path = self.index_dir / "ingested_meta.json"
        self._meta: Dict[str, Any] = {"rows": {}} ## this is dict of rows
        
        if self.meta_path.exists():
            try:
                self._meta = json.loads(self.meta_path.read_text(encoding="utf-8")) or {"rows": {}} # load it if alrady there
            except Exception:
                self._meta = {"rows": {}} # init the empty one if dones not exists
        

        self.model_loader = model_loader or ModelLoader()
        self.emb = self.model_loader.load_embeddings()
        self.vs: Optional[FAISS] = None
        
    def _exists(self)-> bool:
        return (self.index_dir / "index.faiss").exists() and (self.index_dir / "index.pkl").exists()
    
    @staticmethod
    def _fingerprint(text: str, md: Dict[str, Any]) -> str:
        src = md.get("source") or md.get("file_path")
        rid = md.get("row_id")
        if src is not None:
            return f"{src}::{'' if rid is None else rid}"
        return hashlib.sha256(text.encode("utf-8")).hexdigest()
    
    def _save_meta(self):
        self.meta_path.write_text(json.dumps(self._meta, ensure_ascii=False, indent=2), encoding="utf-8")
        
        
    def add_documents(self,docs: List[Document]):
        
        if self.vs is None:
            raise RuntimeError("Call load_or_create() before add_documents_idempotent().")
        
        new_docs: List[Document] = []
        
        for d in docs:
            
            key = self._fingerprint(d.page_content, d.metadata or {})
            if key in self._meta["rows"]:
                continue
            self._meta["rows"][key] = True
            new_docs.append(d)
            
        if new_docs:
            self.vs.add_documents(new_docs)
            self.vs.save_local(str(self.index_dir))
            self._save_meta()
        return len(new_docs)
    
    def load_or_create(self,texts:Optional[List[str]]=None, metadatas: Optional[List[dict]] = None):
        ## if we running first time then it will not go in this block
        if self._exists():
            self.vs = FAISS.load_local(
                str(self.index_dir),
                embeddings=self.emb,
                allow_dangerous_deserialization=True,
            )
            return self.vs
        
        
        if not texts:
            raise DocumentPortalException("No existing FAISS index and no data to create one", sys)
        self.vs = FAISS.from_texts(texts=texts, embedding=self.emb, metadatas=metadatas or [])
        self.vs.save_local(str(self.index_dir))
        return self.vs
        
        
class ChatIngestor:
    def __init__( self,
        temp_base: str = "data",
        faiss_base: str = "faiss_index",
        use_session_dirs: bool = True,
        session_id: Optional[str] = None,
    ):
        try:
            self.model_loader = ModelLoader()
            
            self.use_session = use_session_dirs
            self.session_id = session_id or generate_session_id()
            
            self.temp_base = Path(temp_base); self.temp_base.mkdir(parents=True, exist_ok=True)
            self.faiss_base = Path(faiss_base); self.faiss_base.mkdir(parents=True, exist_ok=True)
            
            self.temp_dir = self._resolve_dir(self.temp_base)
            self.faiss_dir = self._resolve_dir(self.faiss_base)

            log.info("ChatIngestor initialized",
                      session_id=self.session_id,
                      temp_dir=str(self.temp_dir),
                      faiss_dir=str(self.faiss_dir),
                      sessionized=self.use_session)
        except Exception as e:
            log.error("Failed to initialize ChatIngestor", error=str(e))
            raise DocumentPortalException("Initialization error in ChatIngestor", e) from e
            
        
    def _resolve_dir(self, base: Path):
        if self.use_session:
            d = base / self.session_id # e.g. "faiss_index/abc123"
            d.mkdir(parents=True, exist_ok=True) # creates dir if not exists
            return d
        return base # fallback: "faiss_index/"
        
    def _split(self, docs: List[Document], chunk_size=1000, chunk_overlap=200) -> List[Document]:
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        chunks = splitter.split_documents(docs)
        log.info("Documents split", chunks=len(chunks), chunk_size=chunk_size, overlap=chunk_overlap)
        return chunks
    
    def built_retriver( self,
        uploaded_files: Iterable,
        sql_config: Optional[Dict[str, Any]] = None,
        *,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        k: int = 5,):
       
        docs: List[Document] = []

        if uploaded_files and sql_config:
            raise ValueError("Provide either uploaded_files or sql_config, not both.")

        # Load uploaded files
        if uploaded_files:
            paths = save_uploaded_files(uploaded_files, self.temp_dir)
            for p in paths:
                try:
                    if p.suffix.lower() in {".xlsx", ".csv"}:
                        docs.extend(load_excel_or_csv_as_documents(p))
                    else:
                        docs.extend(load_documents([p]))
                except Exception as e:
                    log.error("Failed to process file", file=str(p), error=str(e))
        elif sql_config:
            docs.extend(load_documents_from_sql(
            connection_string=sql_config["connection_string"],
            query=sql_config["query"],
            metadata_fields=sql_config.get("metadata_fields")
        ))

        else:
            raise ValueError("No source provided. Must provide either uploaded_files or sql_config.")

        if not docs:
            raise ValueError("No valid documents loaded")

        chunks = self._split(docs, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        fm = FaissManager(self.faiss_dir, self.model_loader)
        texts = [c.page_content for c in chunks]
        metas = [c.metadata for c in chunks]

        try:
            vs = fm.load_or_create(texts=texts, metadatas=metas)
            added = fm.add_documents(chunks)
            log.info("FAISS index updated", added=added, index=str(self.faiss_dir))
        
        except Exception as e:
             log.error("Failed to build retriever", error=str(e))
             raise DocumentPortalException("Failed to build retriever", e) from e


            
        
            
class DocHandler:
    """
    PDF save + read (page-wise) for analysis.
    """
    def __init__(self, data_dir: Optional[str] = None, session_id: Optional[str] = None):
        self.data_dir = data_dir or os.getenv("DATA_STORAGE_PATH", os.path.join(os.getcwd(), "data", "document_analysis"))
        self.session_id = session_id or generate_session_id("session")
        self.session_path = os.path.join(self.data_dir, self.session_id)
        os.makedirs(self.session_path, exist_ok=True)
        log.info("DocHandler initialized", session_id=self.session_id, session_path=self.session_path)

    def save_file(self, uploaded_file) -> str:
        try:
            filename = os.path.basename(uploaded_file.name)
            ext = Path(filename).suffix.lower()
            if ext not in SUPPORTED_EXTENSIONS:
                raise ValueError(f"Unsupported file type: {ext}. Supported: {SUPPORTED_EXTENSIONS}")
            save_path = os.path.join(self.session_path, filename)
            with open(save_path, "wb") as f:
                if hasattr(uploaded_file, "read"):
                    f.write(uploaded_file.read())
                else:
                    f.write(uploaded_file.getbuffer())
            log.info("File saved successfully", file=filename, save_path=save_path, session_id=self.session_id)
            return save_path
        except Exception as e:
            log.error("Failed to save file", error=str(e), session_id=self.session_id)
            raise DocumentPortalException(f"Failed to save file: {str(e)}", e) from e

    def read_file(self, file_path: str) -> str:
        path = Path(file_path)
        ext = path.suffix.lower()
        try:
            if ext == ".pdf":
                text_chunks = []
                with fitz.open(path) as doc:
                    for page_num in range(doc.page_count):
                        page = doc.load_page(page_num)
                        text_chunks.append(f"\n--- Page {page_num + 1} ---\n{page.get_text()}")
                return "\n".join(text_chunks)

            elif ext == ".docx":
                doc = docx.Document(path)
                return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

            elif ext == ".txt":
                return path.read_text(encoding="utf-8", errors="ignore")
            
            elif ext == ".md":
                return path.read_text(encoding="utf-8", errors="ignore")

            elif ext == ".pptx":
                prs = Presentation(path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                    slides.append(f"\n--- Slide {i + 1} ---\n" + "\n".join(texts))
                return "\n".join(slides)

            elif ext in {".xlsx", ".xls"}:
                df = pd.read_excel(path)
                return df.to_csv(index=False)

            elif ext == ".csv":
                df = pd.read_csv(path)
                return df.to_csv(index=False)

            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            log.error("Failed to read file", error=str(e), file_path=file_path, session_id=self.session_id)
            raise DocumentPortalException(f"Could not process file: {file_path}", e) from e
class DocumentComparator:
    """
    Save, read & combine PDFs for comparison with session-based versioning.
    """
    def __init__(self, base_dir: str = "data/document_compare", session_id: Optional[str] = None):
        self.base_dir = Path(base_dir)
        self.session_id = session_id or generate_session_id()
        self.session_path = self.base_dir / self.session_id
        self.session_path.mkdir(parents=True, exist_ok=True)
        log.info("DocumentComparator initialized", session_path=str(self.session_path))

    def save_uploaded_files(self, reference_file, actual_file):
        files = {"reference": reference_file, "actual": actual_file}
        saved_paths = {}
        try:
            for label, fobj in files.items():
                ext = Path(fobj.name).suffix.lower()
                if ext not in SUPPORTED_EXTENSIONS:
                    raise ValueError(f"Unsupported file type: {ext}")
                out_path = self.session_path / fobj.name
                with open(out_path, "wb") as f:
                        if hasattr(fobj, "read"):
                            f.write(fobj.read())
                        else:
                            f.write(fobj.getbuffer())
                saved_paths[label] = out_path
            log.info("Files saved", reference=str(saved_paths["reference"]),
                     actual=str(saved_paths["actual"]), session=self.session_id)
            return saved_paths["reference"], saved_paths["actual"]
            
        except Exception as e:
            log.error("Error saving files", error=str(e), session=self.session_id)
            raise DocumentPortalException("Error saving files", e) from e
    def read_file(self, file_path: Path) -> str:
        handler = DocHandler(session_id=self.session_id)
        return handler.read_file(str(file_path))

    def combine_documents(self) -> str:
        
        try:
            combined_parts = []
            for file in sorted(self.session_path.iterdir()):
                if file.is_file() and file.suffix.lower() in SUPPORTED_EXTENSIONS:
                    content = self.read_file(file)
                    combined_parts.append(f"Document: {file.name}\n{content}")
            combined_text = "\n\n".join(combined_parts)
            log.info("Documents combined", count=len(combined_parts), session=self.session_id)
            return combined_text
        except Exception as e:
            log.error("Error combining documents", error=str(e), session=self.session_id)
            raise DocumentPortalException("Error combining documents", e) from e

    def clean_old_sessions(self, keep_latest: int = 3):
        try:
            sessions = sorted([f for f in self.base_dir.iterdir() if f.is_dir()], reverse=True)
            for folder in sessions[keep_latest:]:
                shutil.rmtree(folder, ignore_errors=True)
                log.info("Old session folder deleted", path=str(folder))
        except Exception as e:
            log.error("Error cleaning old sessions", error=str(e))
            raise DocumentPortalException("Error cleaning old sessions", e) from e

